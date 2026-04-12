"""Build the presentation PowerPoint from the finalized outline."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# -- Theme colors --
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_SLIDE = RGBColor(0x16, 0x21, 0x3E)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_MAGENTA = RGBColor(0xE0, 0x40, 0xFB)
TABLE_HEADER_BG = RGBColor(0x0D, 0x47, 0xA1)
TABLE_ROW_BG = RGBColor(0x1E, 0x2D, 0x4F)
TABLE_ALT_BG = RGBColor(0x25, 0x35, 0x5A)


def set_slide_bg(slide, color=BG_SLIDE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=TEXT_WHITE, bold=False, space_before=Pt(6), font_name="Meiryo"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    return p


def add_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    return table


def style_table(table, header_data, row_data):
    for i, hd in enumerate(header_data):
        cell = table.cell(0, i)
        cell.text = hd
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_WHITE
            p.font.bold = True
            p.font.name = "Meiryo"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG

    for r, row in enumerate(row_data):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = TEXT_WHITE
                p.font.name = "Meiryo"
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_BG if r % 2 == 0 else TABLE_ALT_BG


def add_part_divider(prs, part_num, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, 1, 2.0, 8, 1, f"Part {part_num}", 24, ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.8, 8, 1.5, title, 36, TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, 1, 4.2, 8, 1, subtitle, 18, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    return slide


def add_slide_title(slide, number, title, color=ACCENT_CYAN):
    add_text_box(slide, 0.5, 0.2, 9, 0.6, f"{number}. {title}", 24, color, bold=True)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # =========================================================================
    # Title Slide
    # =========================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG_DARK)
    add_text_box(sl, 0.5, 1.2, 9, 1.5, "LLM は読めるが語れない", 40, TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, 0.5, 2.6, 9, 1, "期待値プロパティが判定スイッチである", 28, ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, 0.5, 4.0, 9, 0.8, "オントロジー PoC 勉強会発表\n2026-04-12", 16, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    # =========================================================================
    # Part 1: 問題提起
    # =========================================================================
    add_part_divider(prs, 1, "問題提起", "LLM にデータを見せると何が起きるか?")

    # -- Slide 1 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 1, "LLM に売上データを見せると何が起きるか?")
    tf = add_text_box(sl, 0.5, 1.0, 4.5, 3.5, "Q: この四半期売上データを分析して。", 16, TEXT_LIGHT)
    add_para(tf, "")
    add_para(tf, "| 四半期 | 売上       |", 14, TEXT_LIGHT)
    add_para(tf, "| Q1    | 5,148,206  |", 14, TEXT_LIGHT)
    add_para(tf, "| Q2    | 6,181,726  |", 14, TEXT_LIGHT)
    add_para(tf, "| Q3    | 5,486,749  |", 14, TEXT_LIGHT)
    add_para(tf, "| Q4    | 11,014,309 |", 14, TEXT_LIGHT)

    tf2 = add_text_box(sl, 5.2, 1.0, 4.5, 2.0, "LLM の回答:", 16, ACCENT_YELLOW, bold=True)
    add_para(tf2, "「Q4 は Q1-Q3 平均の 1.96 倍です。", 16, TEXT_WHITE)
    add_para(tf2, "  Q4 に売上が集中しています。」", 16, TEXT_WHITE)

    tf3 = add_text_box(sl, 5.2, 3.2, 4.5, 1.5, "", 18, ACCENT_CYAN, bold=True)
    add_para(tf3, "で、この 1.96 倍って...", 20, ACCENT_CYAN, bold=True)
    add_para(tf3, "良いの？ 悪いの？ 想定通りなの？", 20, ACCENT_YELLOW, bold=True)

    add_text_box(sl, 0.5, 4.8, 9, 0.5, "→ LLM は数字を読めている。でも「想定通りか」には答えていない。", 14, TEXT_LIGHT)

    # -- Slide 2 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 2, "LLM は「読む」ことはできるが「判定する」ことはできない")
    add_text_box(sl, 0.5, 1.0, 9, 0.8,
                 "LLM は観測値を正しく記述する。しかし「想定からどれだけずれているか」は、\n期待値を教えない限り語れない。",
                 16, TEXT_WHITE, bold=True)

    tbl = add_table(sl, 3, 4, 0.5, 2.0, 9, 1.5)
    style_table(tbl,
                ["モード", "やっていること", "例", "LLM単体で可能?"],
                [["記述モード", "データにある数字を読む", "\"Q4 は平均の 1.96 倍\"", "✅"],
                 ["判定モード", "想定と比較して評価する", "\"想定 30-50% を 2倍も超えている\"", "❌"]])

    tf = add_text_box(sl, 0.5, 3.8, 9, 1.5, "なぜ「判定モード」に入れないのか?", 16, ACCENT_YELLOW, bold=True)
    add_para(tf, "・「この会社では Q4 は 30-50% 増を想定している」は、その会社にしかない情報", 14, TEXT_LIGHT)
    add_para(tf, "・LLM の訓練データには載っていない", 14, TEXT_LIGHT)
    add_para(tf, "・→ 教えてあげないと、比較のしようがない", 14, ACCENT_CYAN, bold=True)

    # -- Slide 3 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 3, "割引 +7% — これは良いの？ 悪いの？")
    add_text_box(sl, 0.5, 1.0, 9, 0.5, "事実: 割引キャンペーン中の注文数は +7% でした。", 18, TEXT_WHITE, bold=True)

    add_text_box(sl, 0.5, 1.8, 4.2, 0.5, "期待値なし（LLM 単体）", 16, ACCENT_YELLOW, bold=True)
    tf = add_text_box(sl, 0.5, 2.3, 4.2, 1.5, "", 15, TEXT_LIGHT)
    add_para(tf, "「キャンペーン中は注文が", 15, TEXT_LIGHT)
    add_para(tf, "  7% 増えています。supported です。」", 15, TEXT_LIGHT)
    add_para(tf, "→ \"増えてるから OK\" で終わり", 14, ACCENT_RED, bold=True)

    add_text_box(sl, 5.3, 1.8, 4.2, 0.5, "期待値あり（15-30% 増を想定）", 16, ACCENT_GREEN, bold=True)
    tf = add_text_box(sl, 5.3, 2.3, 4.2, 1.5, "", 15, TEXT_LIGHT)
    add_para(tf, "「+7% は想定の 15-30% を大きく", 15, TEXT_LIGHT)
    add_para(tf, "  下回っています。キャンペーン", 15, TEXT_LIGHT)
    add_para(tf, "  設計の見直しが必要です。」", 15, TEXT_LIGHT)
    add_para(tf, "→ 判定 + 次のアクションが出る", 14, ACCENT_GREEN, bold=True)

    add_text_box(sl, 0.5, 4.5, 9, 0.8, "同じ数字。違いは「期待値を知っているかどうか」だけ。", 18, ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

    # -- Slide 4 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 4, "「期待値」を渡すと LLM は判定マシンに変わる")

    tbl = add_table(sl, 4, 3, 0.5, 1.0, 9, 2.0)
    style_table(tbl,
                ["", "期待値なし（記述モード）", "期待値あり（判定モード）"],
                [["出力", "「+7% 増えています」", "「想定 15-30% に対し +7%。期待以下」"],
                 ["判定", "supported（増えたからOK）", "contradicted（想定の半分以下）"],
                 ["次のアクション", "なし", "キャンペーン設計見直し"]])

    add_text_box(sl, 0.5, 3.4, 9, 0.6, "期待値プロパティ ＝ LLM の「判定スイッチ」", 24, ACCENT_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
    tf = add_text_box(sl, 1.5, 4.2, 7, 1.2, "", 18, TEXT_WHITE)
    add_para(tf, "スイッチ OFF → LLM は数字を読むだけの「観測マシン」", 18, TEXT_LIGHT)
    add_para(tf, "スイッチ ON  → LLM は想定と比較する「判定マシン」", 18, ACCENT_GREEN, bold=True)

    # =========================================================================
    # Part 2: オントロジー入門
    # =========================================================================
    add_part_divider(prs, 2, "オントロジー入門", "データの意味づけの 3 層目")

    # -- Slide 5 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 5, "データの意味づけは 3 層で考える")

    layers = [
        ("オントロジー層", "\"なぜそうなるか\" （因果ルール・期待値）", ACCENT_MAGENTA, "「Q4 は季節商品の影響で売上 30-50% 増が想定される」"),
        ("セマンティック層", "\"どう計算するか\" （メトリクス定義）", ACCENT_CYAN, "total_revenue = SUM(total_amount)"),
        ("メタデータ層", "\"何があるか\" （テーブル名・カラム名）", ACCENT_GREEN, "fct_orders テーブルに total_amount カラムがある"),
    ]
    for i, (name, desc, color, example) in enumerate(layers):
        y = 1.2 + i * 1.3
        add_text_box(sl, 0.5, y, 4.5, 0.4, name, 20, color, bold=True)
        add_text_box(sl, 0.5, y + 0.4, 4.5, 0.4, desc, 13, TEXT_LIGHT)
        add_text_box(sl, 5.2, y + 0.1, 4.5, 0.6, example, 13, TEXT_WHITE)

    add_text_box(sl, 0.5, 4.8, 9, 0.5, "LLM にメタデータ＋セマンティックを渡せば SQL は書ける。でも「30-50% 増が想定」はオントロジー層にしかない。", 13, ACCENT_YELLOW)

    # -- Slide 6 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 6, "セマンティックレイヤー = how、オントロジー = why")

    tbl = add_table(sl, 3, 3, 0.5, 1.2, 9, 1.2)
    style_table(tbl,
                ["層", "問い", "例"],
                [["セマンティック", "どう計算するか", "total_revenue = SUM(total_amount)"],
                 ["オントロジー", "なぜそうなるか", "季節商品がQ4に集中 → 売上 30-50% 増"]])

    tf = add_text_box(sl, 0.5, 3.0, 9, 2, "", 16, TEXT_LIGHT)
    add_para(tf, "・dbt Semantic Layer はセマンティック層", 16, TEXT_LIGHT)
    add_para(tf, "・本 PoC が扱うのはその上のオントロジー層", 16, TEXT_WHITE, bold=True)
    add_para(tf, "・多くの組織はメタデータ＋セマンティックまでは整備しているが、", 16, TEXT_LIGHT)
    add_para(tf, "  オントロジー層は未整備", 16, ACCENT_YELLOW, bold=True)

    # -- Slide 7 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 7, "因果ルールは 7 つのプロパティで書ける")

    code_lines = [
        (":DiscountOrderRule", TEXT_LIGHT),
        ("    :cause :DiscountCampaign ;", TEXT_LIGHT),
        ("    :effect :OrderVolume ;", TEXT_LIGHT),
        ("    :direction \"increase\" ;", TEXT_LIGHT),
        ("    :expectedMagnitude \"15-30%\" ;", ACCENT_YELLOW),
        ("    :condition \"discount > 10%\" ;", TEXT_LIGHT),
        ("    :comparedTo :NonCampaignPeriod ;", TEXT_LIGHT),
        ("    :description \"割引...\" .", TEXT_LIGHT),
    ]
    comments = [
        "ルール名",
        "原因: 割引キャンペーン",
        "影響先: 注文数",
        "方向: 増える",
        "★ 期待値: 15-30% 増",
        "条件: 割引率10%以上",
        "比較対象: 非キャンペーン期間",
        "説明文",
    ]
    for i, ((line, color), comment) in enumerate(zip(code_lines, comments)):
        y = 1.0 + i * 0.42
        add_text_box(sl, 0.5, y, 5.0, 0.4, line, 13, color, font_name="Menlo")
        add_text_box(sl, 5.6, y, 4.0, 0.4, f"← {comment}", 12, ACCENT_CYAN if i == 4 else TEXT_LIGHT)

    add_text_box(sl, 0.5, 4.6, 9, 0.7,
                 "⚠ この 7 プロパティは本 PoC で独自に設計。Turtle の記法は W3C 標準だが語彙は独自。",
                 11, TEXT_LIGHT)

    # -- Slide 8 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 8, "期待値プロパティが LLM の「判定スイッチ」")

    tbl = add_table(sl, 8, 3, 0.5, 1.0, 9, 3.2)
    style_table(tbl,
                ["プロパティ", "LLM に教えること", "LLM単体で代替可能?"],
                [["cause（原因）", "「何が原因か」", "✅ 一般常識で推測可能"],
                 ["effect（影響先）", "「何に影響するか」", "✅ 一般常識で推測可能"],
                 ["direction（方向）", "「増えるか減るか」", "✅ 一般常識で推測可能"],
                 ["expectedMagnitude（期待値）", "「どれくらい?」", "❌ 組織固有 → 代替不可能"],
                 ["condition（条件）", "「いつ適用するか」", "⚠ 部分的"],
                 ["comparedTo（比較対象）", "「何と比べるか」", "⚠ 部分的"],
                 ["description（説明）", "補足説明", "✅ 効果小"]])
    # Highlight row 4 (expectedMagnitude)
    for c in range(3):
        cell = tbl.cell(4, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x8B, 0x00, 0x00)

    add_text_box(sl, 0.5, 4.5, 9, 0.8,
                 "上 3 つは frontier LLM が自力でカバー。期待値だけが「人間が書かないと LLM にできないこと」",
                 15, ACCENT_YELLOW, bold=True)

    # =========================================================================
    # Part 3: 実験と結果
    # =========================================================================
    add_part_divider(prs, 3, "実験と結果", "184 実験で検証")

    # -- Slide 9 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 9, "実験設計 — 期待値なし vs あり × 8 つの仮説")

    tf = add_text_box(sl, 0.5, 1.0, 9, 3.5, "", 16, TEXT_WHITE)
    add_para(tf, "シンプルな設計:", 18, ACCENT_CYAN, bold=True)
    add_para(tf, "・同じデータ、同じ仮説を LLM に 2 条件で分析させる", 15, TEXT_LIGHT)
    add_para(tf, "  期待値なし: データと SQL 結果だけを渡す", 15, ACCENT_YELLOW)
    add_para(tf, "  期待値あり: オントロジーの期待値プロパティを追加して渡す", 15, ACCENT_GREEN)
    add_para(tf, "・8 つの仮説 × 2 条件 × 各 5 回 = 80 実験", 15, TEXT_WHITE, bold=True)
    add_para(tf, "・採点: 「LLM が期待値と比較する言語を使ったか」を 0/1 で判定", 15, TEXT_LIGHT)
    add_para(tf, "・盲検二重採点: 採点者はどちらの条件の出力か知らない", 15, TEXT_LIGHT)
    add_para(tf, "")
    add_para(tf, "仮説の選び方:", 18, ACCENT_CYAN, bold=True)
    add_para(tf, "・前半 4 つ: 教科書的な EC 仮説（Q4 売上、割引効果、VIP 単価 etc.）", 15, TEXT_LIGHT)
    add_para(tf, "・後半 4 つ: LLM の常識では解けない仮説（組織固有の目標値、常識と逆方向のデータ）", 15, ACCENT_YELLOW)

    # -- Slide 10 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 10, "結果 — 「読めるが語れない」が一貫")

    tbl = add_table(sl, 4, 3, 1.5, 1.0, 7, 2.0)
    style_table(tbl,
                ["", "期待値なし", "期待値あり"],
                [["期待値比較スコア", "0.05", "0.93"],
                 ["p 値（片側 Wilcoxon）", "—", "p < 0.00002"],
                 ["効果量", "—", "0.88"]])
    # Bold the score cells
    for c in [1, 2]:
        for r in [1, 2, 3]:
            for p in tbl.cell(r, c).text_frame.paragraphs:
                p.font.bold = True

    add_text_box(sl, 0.5, 3.3, 9, 0.5,
                 "→ 期待値プロパティの有無だけで認知モードが切り替わる",
                 18, ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

    tf = add_text_box(sl, 0.5, 4.0, 9, 1.2, "", 14, TEXT_LIGHT)
    add_para(tf, "ただし: LLM はデータを読み間違えているわけではない。", 14, TEXT_LIGHT)
    add_para(tf, "方向判定（「増えた/減った」）は期待値なしでも正解率 90% 以上。", 14, TEXT_LIGHT)
    add_para(tf, "読めるが、\"想定と比べてどうか\" を語れないだけ。", 14, ACCENT_YELLOW, bold=True)

    # -- Slide 11 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 11, "東京 Q4 目標 — LLM 自身が「判定不能」と言った")
    add_text_box(sl, 0.5, 0.9, 9, 0.4, "仮説: 「東京の Q4 売上は社内目標を達成したか?」", 16, TEXT_WHITE)

    # 期待値なし
    add_text_box(sl, 0.5, 1.5, 4.3, 0.4, "期待値なし", 16, ACCENT_YELLOW, bold=True)
    add_text_box(sl, 0.5, 1.9, 4.3, 1.8,
                 "「Q4 東京の売上は 270 万です。\nしかし社内目標を知らないため、\n達成したかどうかは\n判定できません。」",
                 14, TEXT_LIGHT)
    add_text_box(sl, 0.5, 3.6, 4.3, 0.4, "→ 判定不能", 18, ACCENT_RED, bold=True)

    # 期待値あり
    add_text_box(sl, 5.3, 1.5, 4.3, 0.4, "期待値あり（目標 300-400 万）", 16, ACCENT_GREEN, bold=True)
    add_text_box(sl, 5.3, 1.9, 4.3, 1.8,
                 "「Q4 東京の売上は 270 万で、\n目標 300-400 万を\n下回っています。\n未達です。」",
                 14, TEXT_WHITE)
    add_text_box(sl, 5.3, 3.6, 4.3, 0.4, "→ 否定（未達）", 18, ACCENT_GREEN, bold=True)

    tf = add_text_box(sl, 0.5, 4.3, 9, 1.0, "", 15, TEXT_WHITE)
    add_para(tf, "→ LLM 自身が「目標を知らなければ判定できない」と認識している", 15, ACCENT_CYAN, bold=True)
    add_para(tf, "→ 期待値は新しい能力を与えるのではなく、既にある能力に\"比較対象\"を与える", 14, TEXT_LIGHT)

    # -- Slide 12 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 12, "VIP 注文回数 — 常識に騙されないが、判定モードには入らない")
    add_text_box(sl, 0.5, 0.9, 9, 0.4, "仮説: 「VIP 顧客は新規顧客より 1 人あたりの注文数が多い」", 16, TEXT_WHITE)

    tbl = add_table(sl, 3, 2, 0.5, 1.4, 4.5, 1.0)
    style_table(tbl,
                ["", "事実"],
                [["LLM の常識", "「VIP はたくさん買うはず」"],
                 ["データの真実", "VIP 11.3 回 ≈ 新規 11.5 回"]])

    add_text_box(sl, 5.3, 1.4, 4.4, 0.3, "期待値なし", 14, ACCENT_YELLOW, bold=True)
    add_text_box(sl, 5.3, 1.7, 4.4, 0.8, "「VIP は新規とほぼ同じ注文数です。\n仮説は支持されません。」", 13, TEXT_LIGHT)

    add_text_box(sl, 5.3, 2.7, 4.4, 0.3, "期待値あり（+30-50% を想定）", 14, ACCENT_GREEN, bold=True)
    add_text_box(sl, 5.3, 3.0, 4.4, 0.8, "「想定は +30-50% ですが、実測はほぼ同等。\n想定から約 40pp 乖離しています。」", 13, TEXT_WHITE)

    tf = add_text_box(sl, 0.5, 4.0, 9, 1.3, "", 14, TEXT_WHITE)
    add_para(tf, "→ 期待値なしでも常識に騙されずデータを正しく読んでいる（5回中5回正答）", 14, ACCENT_CYAN)
    add_para(tf, "→ でも「想定から 40pp 乖離」という定量的判定は、期待値を渡したときだけ", 14, ACCENT_YELLOW, bold=True)
    add_para(tf, "→ LLM が弱いのではない。比較対象（期待値）がないだけ。", 14, TEXT_WHITE, bold=True)

    # -- Slide 13 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 13, "184 実験の信頼性")

    tbl = add_table(sl, 4, 6, 0.3, 1.0, 9.4, 1.8)
    style_table(tbl,
                ["段階", "実験数", "仮説タイプ", "p 値", "効果量", "Cohen's κ"],
                [["探索的実験", "104", "7仮説 × 6段階", "—", "—", "—"],
                 ["確認実験(第1弾)", "40", "教科書的(4仮説)", "0.0015", "0.45", "0.86"],
                 ["確認実験(第2弾)", "40", "常識で解けない(4仮説)", "< 0.00002", "0.88", "0.95"]])
    # Highlight last row
    for c in range(6):
        cell = tbl.cell(3, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0xA1)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True

    tf = add_text_box(sl, 0.5, 3.2, 9, 2.0, "", 14, TEXT_LIGHT)
    add_para(tf, "・第1弾 → 第2弾で効果量がほぼ 2 倍に（0.45 → 0.88）", 14, TEXT_WHITE)
    add_para(tf, "・第1弾で「期待値なしでも解けるケース」が見つかったが、仮説設計の問題だった（第2弾で解消）", 14, TEXT_LIGHT)
    add_para(tf, "・方法論: 事前登録（実験前にルールを git にコミット）、盲検採点、二重採点", 14, TEXT_LIGHT)
    add_para(tf, "")
    add_para(tf, "→ 最も厳しい条件でも「期待値プロパティの効果」は再現された", 15, ACCENT_CYAN, bold=True)

    # =========================================================================
    # Part 4: デモ
    # =========================================================================
    add_part_divider(prs, 4, "ライブデモ", "期待値なし vs あり を実際に見る")

    # -- Slide 14 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 14, "期待値なし vs あり を実際に見てみよう")

    add_text_box(sl, 0.5, 1.0, 9, 0.4, "ターミナルで実行:", 16, ACCENT_CYAN, bold=True)
    add_text_box(sl, 0.5, 1.5, 9, 0.8,
                 "uv run python -m ontoprobe --llm-demo H7",
                 16, ACCENT_GREEN, font_name="Menlo")

    tf = add_text_box(sl, 0.5, 2.5, 9, 2.5, "", 15, TEXT_LIGHT)
    add_para(tf, "期待される出力:", 16, ACCENT_YELLOW, bold=True)
    add_para(tf, "")
    add_para(tf, "期待値なし:", 14, ACCENT_YELLOW)
    add_para(tf, "  「VIP 11.3 回 ≈ New 11.5 回。差はわずか1%...判定不能」", 14, TEXT_LIGHT)
    add_para(tf, "期待値あり:", 14, ACCENT_GREEN)
    add_para(tf, "  「想定では 30-50% の優位性が期待されるが、実測はほぼ同等。想定と大きく乖離...否定」", 14, TEXT_WHITE)
    add_para(tf, "")
    add_para(tf, "注目ポイント: 出力の語彙が変わる瞬間。「差はわずか」→「想定と大きく乖離」", 15, ACCENT_CYAN, bold=True)

    add_text_box(sl, 0.5, 5.0, 9, 0.4, "⚠ ネットワーク必須。所要 20-30 秒。失敗時はスクショで代替。", 11, TEXT_LIGHT)

    # -- Slide 15 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 15, "Turtle → LLM → 判定 の流れ")

    boxes = [
        (0.8, "Turtle ファイル", ":expectedMagnitude\n\"30-50%\"", ACCENT_MAGENTA),
        (4.0, "LLM プロンプト", "\"想定 30-50%\"\n+ query 結果", ACCENT_CYAN),
        (7.2, "LLM 出力", "\"想定を 2 倍\n超えている\"", ACCENT_GREEN),
    ]
    for x, title, content, color in boxes:
        add_text_box(sl, x, 1.5, 2.5, 0.4, title, 14, color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(sl, x, 2.0, 2.5, 1.2, content, 16, TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    # Arrows
    add_text_box(sl, 3.3, 2.2, 0.7, 0.5, "→", 30, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, 6.5, 2.2, 0.7, 0.5, "→", 30, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    tf = add_text_box(sl, 0.5, 3.8, 9, 1.5, "", 16, TEXT_WHITE)
    add_para(tf, "ポイント:", 16, ACCENT_YELLOW, bold=True)
    add_para(tf, "Turtle の 1 行（expectedMagnitude）が LLM の判定モードを ON にする。", 16, TEXT_WHITE, bold=True)
    add_para(tf, "それ以外のプロパティ（クラス階層、因果方向）は frontier LLM が自力でカバーする。", 14, TEXT_LIGHT)

    # =========================================================================
    # Part 5: 実務への示唆
    # =========================================================================
    add_part_divider(prs, 5, "実務への示唆", "明日から何をすべきか")

    # -- Slide 16 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 16, "期待値の書き方 5 原則")

    tbl = add_table(sl, 6, 4, 0.3, 1.0, 9.4, 2.8)
    style_table(tbl,
                ["原則", "❌ 悪い例", "✅ 良い例", "理由"],
                [["1. 数値範囲で書く", "positive correlation", "15-30%", "曖昧だとLLMも判定不能"],
                 ["2. 条件はSQLで書く", "冬に", "order_quarter = 4", "LLMがSQLに変換可能"],
                 ["3. 比較対象を明示", "（省略）", ":NewCustomer", "ベースライン明確"],
                 ["4. 構造化が先", "長文説明だけ", "構造 + 1行説明", "LLMは構造を優先使用"],
                 ["5. 原因はクラスで指す", "\"vip\" (文字列)", ":VIPCustomer", "テーブル横断可能"]])

    add_text_box(sl, 0.5, 4.2, 9, 1.0,
                 "最大の難所: 「組織の期待値を数値化すること」自体が難しい。\n「なんとなく相関がある」→ 数値範囲にするには、過去データ分析や専門家ヒアリングが必要。",
                 13, TEXT_LIGHT)

    # -- Slide 17 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 17, "どのドメインで期待値投資すべきか")

    tbl = add_table(sl, 4, 3, 1.0, 1.0, 8, 2.0)
    style_table(tbl,
                ["", "LLM が詳しいドメイン", "LLM が弱いドメイン"],
                [["例", "EC、Web マーケ", "半導体、医療、保険"],
                 ["期待値なしの実力", "そこそこ当たる（常識あり）", "当てずっぽう or 自信を持って誤答"],
                 ["期待値の価値", "あると便利", "ないと危険"]])
    # Highlight last row
    for c in range(3):
        cell = tbl.cell(3, c)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = ACCENT_YELLOW

    add_text_box(sl, 0.5, 3.5, 9, 0.6,
                 "→ 自社のドメインが「LLM が詳しくない」領域なら、期待値プロパティの投資は必須",
                 16, ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, 0.5, 4.3, 9, 0.5,
                 "⚠ 注: EC 以外のドメインは本 PoC では実測していません（Masking 実験で用語置換のみ）。上記は推定です。",
                 12, TEXT_LIGHT)

    # -- Slide 18 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl)
    add_slide_title(sl, 18, "限界と次のステップ")

    tf = add_text_box(sl, 0.5, 1.0, 4.5, 4.0, "", 14, TEXT_LIGHT)
    add_para(tf, "本 PoC の限界（正直に）:", 16, ACCENT_RED, bold=True)
    add_para(tf, "")
    add_para(tf, "・Claude モデルのみで検証", 14, TEXT_LIGHT)
    add_para(tf, "  定性的主張は model に依存しないが", 12, TEXT_LIGHT)
    add_para(tf, "  定量値は Claude 固有", 12, TEXT_LIGHT)
    add_para(tf, "・合成データ", 14, TEXT_LIGHT)
    add_para(tf, "  実組織の KPI 目標での検証は未実施", 12, TEXT_LIGHT)
    add_para(tf, "・EC ドメインのみ", 14, TEXT_LIGHT)
    add_para(tf, "  真のドメイン転移は未検証", 12, TEXT_LIGHT)
    add_para(tf, "・8 仮説の代表性は未検証", 14, TEXT_LIGHT)

    tf2 = add_text_box(sl, 5.3, 1.0, 4.3, 4.0, "", 14, TEXT_LIGHT)
    add_para(tf2, "次のステップ:", 16, ACCENT_GREEN, bold=True)
    add_para(tf2, "")
    add_para(tf2, "・他モデルでの再現", 14, TEXT_WHITE)
    add_para(tf2, "  GPT-5 / Gemini 3 / Sonnet", 12, TEXT_LIGHT)
    add_para(tf2, "・実データ・実 KPI 目標での検証", 14, TEXT_WHITE)
    add_para(tf2, "・Snowflake / Cortex への移植", 14, TEXT_WHITE)
    add_para(tf2, "  production 適用", 12, TEXT_LIGHT)

    # =========================================================================
    # Part 6: まとめ
    # =========================================================================
    add_part_divider(prs, 6, "まとめ", "")

    # -- Slide 19 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG_DARK)
    add_text_box(sl, 0.5, 0.8, 9, 1.0,
                 "LLM は読めるが語れない。", 36, TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, 0.5, 1.8, 9, 0.8,
                 "期待値プロパティは、LLM を「観測マシン」から\n「判定マシン」に変えるスイッチである。",
                 24, ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

    tf = add_text_box(sl, 1.0, 3.2, 8, 2.0, "", 16, TEXT_LIGHT)
    add_para(tf, "・LLM は数字を正しく読む（記述モード）", 16, TEXT_LIGHT)
    add_para(tf, "・でも「想定からどれだけずれているか」は自発的に語らない", 16, TEXT_LIGHT)
    add_para(tf, "・期待値を 1 行書くだけで、LLM の認知モードが切り替わる", 16, TEXT_WHITE, bold=True)
    add_para(tf, "・この「想定」は組織にしかない情報。モデルが進化しても外部供給は永久に必要", 16, ACCENT_YELLOW)

    # -- Slide 20 --
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, BG_DARK)
    add_text_box(sl, 0.5, 0.5, 9, 0.8, "Q&A", 40, TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    qa = [
        ("モデルが賢くなれば不要になるのでは?", "期待値は private info。モデル進化では解決しない"),
        ("Opus だけの結果では?", "Haiku でも再現。他モデルは future work（参考資料 A7）"),
        ("EC 以外のドメインでは?", "Masking で一部確認。真のドメイン転移は未実施（参考資料 A8）"),
        ("期待値は誰が書くの?", "ドメイン専門家。最大の難所は曖昧な知識の数値化（参考資料 A10）"),
        ("ルール数が増えたら?", "数百になると SHACL 等の整合性検証が必要"),
    ]
    for i, (q, a) in enumerate(qa):
        y = 1.5 + i * 0.75
        add_text_box(sl, 0.8, y, 4.0, 0.35, f"Q: {q}", 13, ACCENT_CYAN, bold=True)
        add_text_box(sl, 4.9, y, 4.8, 0.35, f"→ {a}", 12, TEXT_LIGHT)

    # =========================================================================
    # Save
    # =========================================================================
    out = "slides/ontoprobe_presentation.pptx"
    import os
    os.makedirs("slides", exist_ok=True)
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
